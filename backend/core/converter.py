import os
from pathlib import Path
from typing import List, Optional, Union
import fitz  # PyMuPDF
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches

from .ai_providers.gemini import GeminiProvider
from .ai_providers.openai import OpenAIProvider
from .ai_providers.anthropic import AnthropicProvider
from .ai_providers.grok import GrokProvider

class SaaSConverter:
    """
    SaaS version of the converter using PyMuPDF for better portability.
    No system dependencies like Poppler required.
    """
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(5.625)

    def __init__(
        self, 
        provider: str = 'gemini',
        api_key: Optional[str] = None,
        model: Optional[str] = None,
        dpi: int = 144, 
        remove_watermark: bool = True
    ):
        self.dpi = dpi
        self.remove_watermark = remove_watermark
        self.provider_name = provider.lower()
        
        # Initialize AI Provider
        if api_key:
            if self.provider_name == 'gemini':
                self.ai_provider = GeminiProvider(api_key, model or "gemini-2.5-flash")
            elif self.provider_name == 'openai':
                self.ai_provider = OpenAIProvider(api_key, model or "gpt-4o")
            elif self.provider_name in ['anthropic', 'claude']:
                self.ai_provider = AnthropicProvider(api_key, model or "claude-3-5-sonnet-20241022")
            elif self.provider_name in ['grok', 'xai']:
                # grok-2-1212 confirmed to support vision
                self.ai_provider = GrokProvider(api_key, model or "grok-2-1212")
            else:
                raise ValueError(f"지원하지 않는 AI 프로바이더입니다: {provider}")
        else:
            self.ai_provider = None

    def convert_pdf_to_images(self, pdf_path: Union[str, Path]) -> List[Image.Image]:
        """Convert PDF to images using PyMuPDF."""
        doc = fitz.open(pdf_path)
        images = []
        
        for i in range(len(doc)):
            page = doc.load_page(i)
            # Render page to a pixmap
            pix = page.get_pixmap(matrix=fitz.Matrix(self.dpi/72, self.dpi/72))
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            if self.remove_watermark:
                img = self._remove_watermark(img)
                
            images.append(img)
            
        doc.close()
        return images

    def _remove_watermark(self, image: Image.Image) -> Image.Image:
        """Remove NotebookLM watermark by masking bottom-right area."""
        from PIL import ImageDraw
        width, height = image.size
        m_w = int(width * 0.18)
        m_h = int(height * 0.08)
        
        x0, y0 = width - m_w, height - m_h
        x1, y1 = width, height
        
        # Sample background color
        sample_x = max(0, x0 - 5)
        sample_y = max(0, y0 - 5)
        bg_color = image.getpixel((sample_x, sample_y))
        
        draw = ImageDraw.Draw(image)
        draw.rectangle([x0, y0, x1, y1], fill=bg_color)
        return image

    def create_pptx(
        self, 
        images: List[Image.Image], 
        output_path: Union[str, Path],
        generate_notes: bool = True,
        context: Optional[str] = None
    ) -> Path:
        """Create PPTX from images and generate notes."""
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        blank_layout = prs.slide_layouts[6]

        for idx, img in enumerate(images):
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image
            with io.BytesIO() as img_buffer:
                img.save(img_buffer, format="PNG")
                img_buffer.seek(0)
                slide.shapes.add_picture(
                    img_buffer,
                    Inches(0), Inches(0),
                    width=self.SLIDE_WIDTH, height=self.SLIDE_HEIGHT
                )
            
            # Generate notes
            if generate_notes and self.ai_provider:
                try:
                    notes = self.ai_provider.analyze_slide(img, context)
                    notes_slide = slide.notes_slide
                    notes_frame = notes_slide.notes_text_frame
                    notes_frame.text = notes
                except Exception as e:
                    print(f"Notes generation failed for slide {idx+1}: {e}")
        
        prs.save(str(output_path))
        return Path(output_path)

    def convert(
        self,
        pdf_path: Union[str, Path],
        output_path: Union[str, Path],
        generate_notes: bool = True,
        context: Optional[str] = None
    ) -> Path:
        """Full conversion pipeline."""
        images = self.convert_pdf_to_images(pdf_path)
        return self.create_pptx(images, output_path, generate_notes, context)
